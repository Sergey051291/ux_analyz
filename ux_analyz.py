# -*- coding: utf-8 -*-  # Указываем кодировку файла (нужно для корректной работы с кириллицей)

import time, base64, uuid, os  # time — задержки/тайминги; base64 — кодирование; uuid — уникальные идентификаторы; os — работа с файловой системой
from datetime import datetime  # Работа с датой/временем (для меток времени в отчёте)
from functools import lru_cache  # Декоратор для кэширования результатов функций (например, токенов)
from enum import Enum  # Базовый класс для перечислений (Enum)
from typing import Optional, Dict, List  # Подсказки типов для удобства и читаемости
from io import BytesIO  # Буфер в памяти (нужен для PPTX и скриншотов без записи на диск)

import streamlit as st  # Фреймворк для веб-интерфейса
from pptx import Presentation  # Создание/редактирование PowerPoint презентации
from pptx.util import Inches, Pt  # Единицы измерения: дюймы, пункты (размеры и шрифты в PPTX)
from pptx.enum.text import PP_ALIGN  # type: ignore  # Константы выравнивания текста в PPTX (игнорим подсказки типов)
from pptx.enum.shapes import PP_PLACEHOLDER  # type: ignore  # Типы плейсхолдеров на слайдах (игнорим подсказки типов)
from selenium import webdriver  # Драйвер Selenium для управления браузером
from selenium.webdriver.chrome.options import Options  # Настройки Chrome
from selenium.webdriver.chrome.service import Service  # Обёртка сервиса для chromedriver
import requests  # HTTP-клиент для загрузки HTML страниц
import ollama  # Клиент для локального LLM сервера Ollama

# ==================== КОНФИГ ====================

class LLMProvider(Enum):  # Перечисление доступных LLM-провайдеров
    OLLAMA = "ollama"     # Локальная модель через Ollama
    GIGACHAT = "gigachat" # Облачный API GigaChat

CONFIG = {  # Глобальный конфиг приложения
    "llm_provider": LLMProvider.OLLAMA,  # Текущий провайдер по умолчанию — Ollama
    "gigachat": {                        # Параметры доступа к GigaChat
        "client_id": "ВАШ_CLIENT_ID",    # Ваш client_id
        "client_secret": "ВАШ_SECRET_PART",  # Ваш client_secret (секретная часть)
        "scope": "GIGACHAT_API_PERS",    # Скоуп доступа
        "auth_url": "https://ngw.devices.sberbank.ru:9443/api/v2/oauth",  # URL получения токена
        "api_url": "https://gigachat.devices.sberbank.ru/api/v1/chat/completions"  # URL API чата
    },
    "ollama": {"model": "llama3", "host": "http://localhost:11434"},  # Настройки Ollama
    "max_pages": 5,                # Максимум анализируемых страниц за запуск
    "max_html_length": 15000,      # Ограничение длины HTML (сокращаем, чтобы не перегружать модель)
    "analysis_modes": ["Базовый", "Расширенный", "Полный"]  # Режимы анализа (пока информативно)
}

# ==================== LLM КЛИЕНТ ====================

class LLMClient:
    def __init__(self, config: Dict):  # Инициализация клиента с конфигом
        self.config = config           # Сохраняем конфиг
        self.provider = config["llm_provider"]  # Запоминаем текущего провайдера

    @lru_cache(maxsize=1)  # Кэшируем токен GigaChat на время работы (чтобы не запрашивать каждый раз)
    def _get_gigachat_token(self) -> Optional[str]:  # Получение access_token для GigaChat
        try:
            auth_b64 = base64.b64encode(  # Формируем Base64 из client_id:client_secret
                f"{self.config['gigachat']['client_id']}:{self.config['gigachat']['client_secret']}".encode()
            ).decode("utf-8")
            headers = {  # Заголовки запроса для OAuth
                'Authorization': f'Basic {auth_b64}',          # Basic авторизация
                'RqUID': str(uuid.uuid4()),                    # Уникальный ID запроса
                'Content-Type': 'application/x-www-form-urlencoded'  # Тип тела
            }
            r = requests.post(  # Запрашиваем токен
                self.config['gigachat']['auth_url'],
                headers=headers,
                data={'scope': self.config['gigachat']['scope']},  # Передаём scope
                verify=False,   # Отключаем проверку SSL (для локальных/тестовых окружений — на проде включить!)
                timeout=15      # Таймаут запроса
            )
            if r.status_code == 200:        # Если всё ок
                return r.json()['access_token']  # Возвращаем токен
            st.error(f"Ошибка аутентификации GigaChat: {r.text}")  # Иначе показываем ошибку
            return None  # Возвращаем None при ошибке
        except Exception as e:  # Ловим любые исключения
            st.error(f"Ошибка получения токена GigaChat: {e}")  # Показываем ошибку
            return None  # И возвращаем None

    def generate(self, prompt: str) -> Optional[str]:  # Универсальный метод генерации ответа LLM
        if self.provider == LLMProvider.OLLAMA:        # Если выбран Ollama
            try:
                resp = ollama.generate(                # Отправляем запрос локальной модели
                    model=self.config['ollama']['model'],  # Имя модели
                    prompt=f"Отвечай только на русском языке. {prompt}",  # Промпт с требованием RU
                    system="Ты профессиональный UX-аналитик. Все ответы должны быть строго на русском языке.",  # Системное сообщение
                    stream=False  # Без стриминга (получаем готовый ответ)
                )
                return resp['response']  # Возвращаем текст ответа
            except Exception as e:        # Ошибки Ollama
                st.error(f"Ошибка Ollama: {e}")  # Показываем ошибку
                st.info("Убедитесь, что Ollama запущен: 'ollama serve'")  # Подсказка
                return None  # Возвращаем None
        else:  # Иначе используем GigaChat
            token = self._get_gigachat_token()  # Получаем токен
            if not token:                       # Если не удалось — выходим
                return None
            headers = {'Authorization': f'Bearer {token}', 'Content-Type': 'application/json'}  # Заголовки API
            payload = {  # Тело запроса к GigaChat
                "model": "GigaChat:latest",  # Модель
                "messages": [                 # История сообщений
                    {"role": "system",
                     "content": "Ты профессиональный UX-аналитик. Все ответы должны быть на русском языке."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.7,  # Температура (степень креативности)
                "max_tokens": 3000   # Лимит токенов ответа
            }
            try:
                r = requests.post(self.config['gigachat']['api_url'], headers=headers, json=payload, timeout=30)  # Запрос
                r.raise_for_status()  # Бросаем исключение при HTTP ошибке
                return r.json()['choices'][0]['message']['content']  # Возвращаем текст ответа
            except Exception as e:  # Ловим ошибки запроса
                st.error(f"Ошибка API GigaChat: {e}")  # Показываем ошибку
                return None  # Возвращаем None


llm_client = LLMClient(CONFIG)  # Создаём экземпляр клиента LLM с текущим конфигом

# ==================== УТИЛИТЫ ====================

def setup_selenium():
    """Запускает Chrome с автоподбором драйвера.
    Сначала пробуем Selenium Manager (встроен в Selenium 4.6+),
    если не вышло — резерв через webdriver-manager."""
    opts = Options()                                  # создаём объект настроек Chrome
    opts.add_argument("--headless=new")               # запуск в новом headless-режиме (Chromium 109+), без окна
    opts.add_argument("--disable-gpu")                # отключаем GPU (устраняет артефакты/краши на CI/Windows)
    opts.add_argument("--no-sandbox")                 # нужно в контейнерах/под root; иначе Chrome может не стартовать
    opts.add_argument("--window-size=1366,900")       # фиксируем размер окна для стабильных скриншотов/вёрстки

    # 1) Пытаемся использовать Selenium Manager (автоматически найдёт/скачает подходящий chromedriver)
    try:
        return webdriver.Chrome(options=opts)         # если получилось — сразу возвращаем готовый драйвер
    except Exception as e_primary:                    # запоминаем причину падения первой попытки
        # 2) Резерв: webdriver-manager (скачает конкретный chromedriver под текущий Chrome)
        try:
            from webdriver_manager.chrome import ChromeDriverManager  # импортируем менеджер драйвера "на лету"
            from selenium.webdriver.chrome.service import Service     # Service для явной передачи пути к драйверу
            service = Service(ChromeDriverManager().install())        # скачиваем (если нужно) и берём путь к драйверу
            return webdriver.Chrome(service=service, options=opts)    # запускаем Chrome с этим драйвером
        except Exception as e_fallback:                # если и резервная попытка провалилась —
            # бросаем одно понятное исключение с ДВУМЯ первопричинами — так проще диагностировать
            raise RuntimeError(
                f"Не удалось запустить Chrome.\n"                     # общее описание проблемы
                f"Selenium Manager: {e_primary}\n"                    # что пошло не так в первой попытке
                f"webdriver-manager: {e_fallback}"                    # и что пошло не так в резервной
            )
def get_page_html(url: str):
    """Загружаем HTML: сначала requests, если не получилось — резерв через Selenium.page_source."""
    # Нормализуем схему: если ввели без http/https — добавим https
    u = (url or "").strip()
    if u and not u.startswith(("http://", "https://")):
        u = "https://" + u

    # Попытка №1 — requests (быстро и дёшево)
    try:
        r = requests.get(
            u,
            headers={
                # Браузерный UA чуть реже ловит антибот
                "User-Agent": (
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/124.0 Safari/537.36"
                )
            },
            timeout=20,
            allow_redirects=True,
        )
        r.raise_for_status()
        return r.text[:CONFIG["max_html_length"]]
    except Exception as e_req:
        # Попытка №2 — Selenium (часто помогает против JS-редиректов/антибота)
        d = None
        try:
            d = setup_selenium()
            d.get(u)
            time.sleep(2.5)  # дать странице дорисоваться
            html = d.page_source
            return html[:CONFIG["max_html_length"]]
        except Exception as e_sel:
            st.error(f"Ошибка загрузки страницы {u}: requests: {e_req}; selenium: {e_sel}")
            return None
        finally:
            try:
                if d is not None:
                    d.quit()
            except Exception:
                pass

def take_screenshots(url: str, max_shots: int = 3) -> List[BytesIO]:
    """Делаем до max_shots скриншотов, прокручивая страницу."""
    shots: List[BytesIO] = []  # Будем собирать скриншоты в виде потоков в памяти
    try:
        d = setup_selenium()   # Запускаем Selenium (Chrome)
        d.get(url)             # Открываем нужный URL
        time.sleep(2)          # Даём странице стабилизироваться
        for _ in range(max_shots):                       # Делаем указанное число кадров
            png = d.get_screenshot_as_png()             # Получаем PNG скриншот
            shots.append(BytesIO(png))                  # Кладём в список как BytesIO
            d.execute_script("window.scrollBy(0, window.innerHeight);")  # Прокручиваем на один экран вниз
            time.sleep(1)                               # Пауза, чтобы прогрузился контент
        d.quit()                                        # Закрываем браузер
    except Exception as e:                               # Ловим ошибки Selenium/Chrome
        st.warning(f"⚠️ Скриншоты не получились: {e}")  # Показываем предупреждение (но не валим процесс)
    return shots                                         # Возвращаем (возможно пустой) список скринов


def analyze_page(html: str, page_name: str):  # Подготовка промпта и вызов LLM для анализа UX
    prompt = f"""
Проанализируй UX страницы {page_name} на русском языке:
1. Структура и навигация
2. Конверсионные элементы
3. Выявленные проблемы
4. Рекомендации по улучшению

Контент:
{html[:10000]}... [сокращено]
"""
    return llm_client.generate(prompt)  # Возвращаем ответ модели (или None при ошибке)


def chunk_text(text: str, chunk_len: int = 900) -> List[str]:
    """Режем текст на компактные куски, чтобы наверняка влезал на слайд."""
    text = (text or "").strip()  # Страхуемся от None и обрезаем пробелы по краям
    if not text:                  # Если пусто —
        return []                 # Возвращаем пустой список
    return [text[i:i + chunk_len] for i in range(0, len(text), chunk_len)]  # Режем на фрагменты по chunk_len


# ---------- Надёжная работа с плейсхолдерами в PPTX ----------

def _get_title_shape(slide):  # Ищем заголовок на слайде
    """Возвращает заголовок, если есть; иначе None."""
    if slide.shapes.title is not None:  # Если у слайда явно есть заголовок —
        return slide.shapes.title       # Возвращаем его
    for ph in slide.placeholders:       # Иначе проходим по плейсхолдерам
        try:
            if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):  # Если это заголовок
                return ph  # Возвращаем найденный плейсхолдер
        except Exception:
            pass  # Игнорируем странные/нестандартные плейсхолдеры
    return None  # Если не нашли — возвращаем None


def _get_body_shape(slide):  # Ищем основной текстовый плейсхолдер
    """Ищет плейсхолдер для основного текста; если нет — None."""
    preferred = (  # Набор типов плейсхолдеров, которые подойдут как "тело" слайда
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.OBJECT,
        PP_PLACEHOLDER.TABLE,
        PP_PLACEHOLDER.PICTURE,
        PP_PLACEHOLDER.CHART,
        PP_PLACEHOLDER.SUBTITLE,
    )
    for ph in slide.placeholders:  # Проходим по плейсхолдерам
        try:
            if getattr(ph.placeholder_format, "type", None) in preferred:  # Если тип нам подходит —
                return ph  # Возвращаем
        except Exception:
            pass  # Игнорируем ошибки доступа к атрибутам
    return None  # Если ничего не нашли — None


def _ensure_text_frame(slide, left=Inches(0.8), top=Inches(1.3), width=Inches(11.8), height=Inches(5.2)):
    """
    Гарантированно возвращает text_frame для основного текста:
    - если есть body-плейсхолдер — используем его
    - иначе создаём textbox вручную
    """
    body = _get_body_shape(slide)  # Пробуем найти готовый плейсхолдер
    if body is not None and body.has_text_frame:  # Если есть и у него есть текстовый фрейм —
        return body.text_frame                    # Возвращаем его
    tb = slide.shapes.add_textbox(left, top, width, height)  # Иначе создаём новый textbox
    return tb.text_frame  # Возвращаем его текстовый фрейм


def _set_title(slide, text):  # Унифицированная установка заголовка
    """Надёжно ставит заголовок (32pt). Если нет заголовка — создаёт textbox."""
    title = _get_title_shape(slide)  # Ищем существующий заголовок
    if title is None:  # Если заголовка нет —
        title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(1.0))  # Создаём textbox
    tf = title.text_frame  # Берём текстовый фрейм
    tf.clear()             # Очищаем его
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()  # Гарантируем первый параграф
    p.text = text          # Ставим текст
    p.font.size = Pt(32)   # Размер шрифта заголовка — 32 pt
    return title           # Возвращаем shape заголовка


def create_ux_presentation_bytes(
    analyses: dict,                         # Словарь {имя страницы: текст анализа}
    screenshots: Dict[str, List[BytesIO]],  # Словарь {имя страницы: список скриншотов BytesIO}
    min_slides: int = 3,                    # Минимум слайдов на сайт (комбинированно: текст + скриншоты)
    max_slides: int = 5                     # Максимум слайдов на сайт
) -> BytesIO:
    """
    На каждый сайт: 3–5 слайдов.
      • 1–2 текстовых (компактная верстка)
      • Остальное — слайды со скриншотами (по одному на слайд)
    Заголовки 32 pt, текст плотно под заголовком.
    """
    prs = Presentation()                 # Создаём презентацию
    prs.slide_width = Inches(13.33)      # Ширина слайда (16:9)
    prs.slide_height = Inches(7.5)       # Высота слайда (16:9)

    # титульный
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Добавляем титульный слайд (layout 0)
    _set_title(slide, "Отчет по UX анализу")            # Устанавливаем заголовок
    try:
        slide.placeholders[1].text = f"Сгенерировано: {datetime.now().strftime('%d.%m.%Y %H:%M')}"  # Подзаголовок/дата
    except Exception:
        cap = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.8), Inches(0.8))  # Если плейсхолдера нет — textbox
        cap.text_frame.text = f"Сгенерировано: {datetime.now().strftime('%d.%m.%Y %H:%M')}"  # Пишем в textbox

    title_only_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]  # Лэйаут с титулом
    text_layout = prs.slide_layouts[1]  # Лэйаут "Title and Content" чаще всего индекс 1

    screenshots = screenshots or {}  # Страхуемся от None

    def _apply_compact_paragraph(p):  # Функция компактного форматирования параграфа
        p.font.size = Pt(18)          # Шрифт текста — 18 pt
        p.alignment = PP_ALIGN.LEFT   # Выравнивание по левому краю
        p.space_before = Pt(0)        # Без отступа перед абзацем
        p.space_after = Pt(0)         # Без отступа после абзаца
        p.line_spacing = 1.0          # Межстрочный интервал 1.0 (плотно)

    for name, analysis in analyses.items():  # Идём по всем анализам
        text_chunks = chunk_text(analysis, chunk_len=900)  # Режем текст на куски по ~900 символов
        shots = [s for s in screenshots.get(name, []) if s is not None]  # Берём скриншоты для страницы (фильтруем None)

        text_slides = min(2, max(1, len(text_chunks)))  # Кол-во текстовых слайдов: 1 или 2 (если кусков мало/много)
        available_imgs = len(shots)                     # Сколько есть скриншотов
        img_slots = min(available_imgs, max_slides - text_slides)  # Сколько сможем показать скриншотов в лимите

        if text_slides + img_slots < min_slides and len(text_chunks) > text_slides:  # Если не дотягиваем до минимума
            text_slides = min(2, text_slides + (min_slides - (text_slides + img_slots)))  # Добавим ещё один текстовый
            img_slots = min(available_imgs, max_slides - text_slides)  # Пересчитаем слоты под скриншоты

        # --- текстовые ---
        for t_i in range(text_slides):         # Рисуем каждый текстовый слайд
            s = prs.slides.add_slide(text_layout)  # Создаём слайд с текстовым лэйаутом
            _set_title(s, f"Анализ: {name} — часть {t_i+1}" if text_slides > 1 else f"Анализ: {name}")  # Заголовок

            # плотный текстовый блок чуть ближе к заголовку
            body_tf = _ensure_text_frame(   # Получаем или создаём текстовый блок
                s,
                left=Inches(0.8),
                top=Inches(1.1),       # Было 1.3 — пододвинули ближе к заголовку
                width=Inches(11.8),
                height=Inches(5.4)     # Чуть выше, чтобы больше текста влезало
            )
            body_tf.clear()            # Очищаем текстовый фрейм
            body_tf.word_wrap = True   # Перенос слов включён
            # убираем внутренние поля у текст-кадра (если поддерживается)
            try:
                body_tf.margin_top = 0
                body_tf.margin_bottom = 0
                body_tf.margin_left = 0
                body_tf.margin_right = 0
            except Exception:
                pass  # На всякий случай игнорируем, если в версии библиотеки нет этих полей

            p = body_tf.add_paragraph()  # Добавляем один параграф (весь кусок текста кладём в него)
            p.text = text_chunks[t_i] if t_i < len(text_chunks) else ""  # Сам текст
            _apply_compact_paragraph(p)  # Применяем плотное форматирование

        # --- скриншоты по одному на слайд ---
        if img_slots > 0:                            # Если есть что показывать
            for j in range(img_slots):               # Делаем слайды под каждый скриншот
                if j >= len(shots):                  # Перестраховка: если внезапно слотов больше чем скринов
                    break
                s = prs.slides.add_slide(title_only_layout)  # Слайд с заголовком
                _set_title(s, f"Анализ: {name} — скриншот {j+1}")  # Заголовок с номером скрина
                try:
                    shots[j].seek(0)  # Прокручиваем поток в начало (на случай, если кто-то уже читал его)
                except Exception:
                    pass
                left = Inches(0.8)   # Левый отступ картинки
                top = Inches(1.4)    # Верхний отступ (под заголовком)
                width = Inches(11.8) # Ширина картинки (автоподбор высоты)
                try:
                    s.shapes.add_picture(shots[j], left, top, width=width)  # Вставляем изображение
                except Exception as e:  # Ловим ошибки вставки изображений
                    print(f"⚠️ Не удалось вставить скриншот {j+1} для «{name}»: {e}")  # Пишем в консоль (не валим процесс)

    bio = BytesIO()  # Создаём буфер в памяти под итоговый PPTX
    prs.save(bio)    # Сохраняем презентацию в буфер
    bio.seek(0)      # Возвращаемся в начало буфера
    return bio       # Возвращаем BytesIO с презентацией


# ==================== ИНТЕРФЕЙС ====================

st.set_page_config(layout="wide", page_title="UX Анализатор", page_icon="🔍")  # Настройки страницы Streamlit

# Спрятать статус как можно раньше, чтобы не было «вспышки»
st.markdown(
    "<style>header [data-testid='stStatusWidget']{display:none!important}</style>",
    unsafe_allow_html=True,
)

# состояние
if "is_running" not in st.session_state: st.session_state.is_running = False  # Флаг: идёт ли анализ
if "cancel" not in st.session_state:     st.session_state.cancel = False      # Флаг: запрос на отмену
if "analyses_last" not in st.session_state: st.session_state.analyses_last = {}  # Последние результаты анализа
if "urls_last" not in st.session_state:     st.session_state.urls_last = {}      # Последние URL'ы
if "shots_last" not in st.session_state:    st.session_state.shots_last = {}     # Последние скриншоты

st.markdown("""
<style>                                                /* открываем тег со стилями */

/* ---- базовая переменная для вертикального сдвига шапки ---- */
:root { --st-header-h: -10px; }                       /* объявляем CSS-переменную: отриц. значение поднимает контент вверх */

/* ---- основной контейнер контента (разные варианты селекторов под версии Streamlit) ---- */
div[data-testid="stAppViewContainer"] > .main .block-container,  /* первый вариант пути к блоку */
main[data-testid="stAppViewContainer"] .block-container,         /* альтернативный вариант (через main) */
div[data-testid="stAppViewContainer"] .block-container{          /* запасной вариант (короче) */
  margin-top: calc(-1 * var(--st-header-h)) !important;          /* сдвигаем контейнер вверх на величину переменной */
  padding-top: .5rem !important;                                  /* и даём небольшой внутренний отступ сверху */
}

/* ---- контейнер сайдбара (тоже в нескольких вариантах, как выше) ---- */
aside[data-testid="stSidebar"] .block-container,                 /* aside-версия селектора */
section[data-testid="stSidebar"] .block-container,               /* section-версия селектора */
div[data-testid="stSidebar"] .block-container{                   /* div-версия селектора */
  margin-top: calc(-1 * var(--st-header-h)) !important;          /* поднимаем содержимое сайдбара */
  padding-top: .5rem !important;                                  /* и добавляем небольшой верхний отступ */
}

/* ---- уплотняем вертикальные отступы заголовков ---- */
h1, h2, h3 { margin-top: .25rem !important; }                    /* уменьшаем верхний margin у h1–h3 */

/* ---- скрываем статус-виджет в хедере (зелёная/жёлтая точка и т. п.) ---- */
header [data-testid="stStatusWidget"] { display: none !important; }  /* полностью прячем виджет */

/* ---- делаем кнопки компактнее (ширина по содержимому, меньше внешних отступов) ---- */
div.stButton > button {
  width: auto !important;                                        /* ширина по контенту, не на всю колонку */
  margin: 0 !important;                                          /* убираем внешние отступы */
  padding: 0.4rem 0.9rem !important;                             /* компактные внутренние отступы */
}

/* ---- уменьшаем промежуток между кнопками в первой колонке формы ---- */
[data-testid="stForm"] [data-testid="stHorizontalBlock"] > div:first-child [data-testid="stHorizontalBlock"]{
  gap: 6px !important;                                           /* расстояние между элементами — 6px */
}

/* ---- убираем «ссылки-якоря» у заголовков (иконки-ссылки справа) ---- */
h1 a[href^="#"],                                                 /* ссылка внутри h1, начинающаяся с # */
h2 a[href^="#"],                                                 /* то же для h2 */
h3 a[href^="#"],                                                 /* то же для h3 */
h4 a[href^="#"],                                                 /* то же для h4 */
a.stHeading-link,                                                /* устаревший/альтернативный класс ссылки */
a[data-testid="stHeadingAnchor"] {                               /* актуальный якорь заголовка */
  display: none !important;                                      /* полностью скрываем элемент */
  visibility: hidden !important;                                 /* дополнительно скрываем видимость */
  pointer-events: none !important;                               /* отключаем клики */
  opacity: 0 !important;                                         /* делаем прозрачным (на всякий случай) */
}

/* =========================
   одинаковая высота «серых» полей под URL
   ========================= */

/* одна ручка-переменная для обеих высот; меняй число — изменятся оба поля */
:root{ --url-box-h: 96px; }                                      /* комфортный диапазон: 88–110px */

/* ---- левое поле: st.text_input (увеличиваем СЕРУЮ ОБЁРТКУ, а не сам <input>) ---- */
div[data-testid="stTextInput"] div[data-baseweb="input"]{
  min-height: var(--url-box-h) !important;                       /* растягиваем серую «рамку» до нужной высоты */
  display: flex !important;                                      /* делаем flex-контейнер */
  align-items: flex-start !important;                             /* выравниваем содержимое по верхнему краю */
  padding-top: 6px !important;                                    /* небольшой отступ сверху, чтобы текст не прилипал */
}

/* сам <input> оставляем нормальным — чтобы placeholder не «расползался» и был у верхней границы */
div[data-testid="stTextInput"] input[aria-label="URL вашего сайта"]{
  height: auto !important;                                       /* высота по содержимому */
  line-height: normal !important;                                 /* стандартная межстрочная высота */
  padding-top: .5rem !important;                                  /* аккуратные внутренние отступы */
  padding-bottom: .5rem !important;                               /* снизу — тоже */
}

/* ---- правое поле: st.text_area ---- */
/* тянем серую обёртку вокруг textarea до той же высоты */
div[data-testid="stTextArea"] div[data-baseweb="textarea"]{
  min-height: var(--url-box-h) !important;                        /* рост контейнера до нужной высоты */
}
/* фиксируем высоту самой textarea (учитываем два возможных текста подписи в aria-label) */
div[data-testid="stTextArea"] textarea[aria-label="URL конкурентов (по одному на строку)"],
div[data-testid="stTextArea"] textarea[aria-label="URL-адрес аналогичный (по одному слову)"]{
  height: var(--url-box-h) !important;                            /* высота поля */
  min-height: var(--url-box-h) !important;                        /* не давать сжиматься меньше */
  max-height: var(--url-box-h) !important;                        /* не давать расти больше */
}
</style>                                                          /* закрываем тег стилей */
""", unsafe_allow_html=True)                                      # разрешаем вставку сырого HTML/CSS


st.title("🔍 Профессиональный UX-анализатор")  # Заголовок страницы

# Сайдбар
with st.sidebar:  # Боковая панель настроек
    st.header("🔧 Настройки")  # Заголовок секции
    provider = st.selectbox("LLM провайдер", [p.value for p in LLMProvider],  # Выбор провайдера LLM
                            index=0 if CONFIG['llm_provider'] == LLMProvider.OLLAMA else 1)
    CONFIG['llm_provider'] = LLMProvider(provider)  # Сохраняем выбранного провайдера
    llm_client.provider = CONFIG['llm_provider']    # Обновляем в клиенте LLM
    if st.button("Проверить подключение"):
        try:
            if llm_client.provider == LLMProvider.GIGACHAT:
                token = llm_client._get_gigachat_token()
                if token:
                    st.success("✅ Успешное подключение к GigaChat!")
                # если токена нет — _get_gigachat_token уже показал st.error
            else:
                ollama.list()
                st.success("✅ Ollama работает корректно!")
                st.code(f"Используется модель:\n{CONFIG['ollama']['model']}")
        except Exception as e:
            st.error(f"Ошибка при проверке подключения: {e}")
    if st.button("Обновить токен") and llm_client.provider == LLMProvider.GIGACHAT:  # Сброс кэша токена GigaChat
        llm_client._get_gigachat_token.cache_clear()  # Очищаем кэш
        st.success("Кэш токена очищен!")              # Сообщаем
    st.markdown("---")                                 # Разделитель
    st.header("⚙️ Параметры анализа")                  # Заголовок секции параметров
    analysis_mode = st.selectbox("Режим анализа", CONFIG['analysis_modes'], index=1)  # Режим анализа (информативно)
    include_screenshots = st.checkbox("Делать скриншоты", True)                        # Флаг: делать скрины или нет

# Форма
with st.form("analysis_form"):     # Начинаем форму (важно для работы двух кнопок в одной форме)
    col1, col2 = st.columns([2, 3])  # Две колонки: слева сайт, справа конкуренты

    with col1:  # Левая колонка
        hcol, linkcol = st.columns([12, 1])  # Внутри две колонки — заголовок и маленькая колонка под ссылку-якорь
        with hcol:
            st.subheader("📌 Анализируемый сайт")  # Подзаголовок секции
        with linkcol:
            st.markdown('<div style="text-align:right; padding-top:10px;"><a href="#exp-main">🔗</a></div>',
                        unsafe_allow_html=True)  # Ссылка-якорь к нижнему итоговому блоку

        main_url = st.text_input("URL вашего сайта", key="main_url",  # Поле ввода основного URL
                                 placeholder="https://ваш-сайт.ru", help="Основной сайт для анализа UX")

        b1, b2 = st.columns([1, 1])  # Две кнопки в одной строке
        with b1:
            submitted = st.form_submit_button("🚀 Начать анализ", type="primary")  # Кнопка запуска анализа
        with b2:
            stop_clicked = st.form_submit_button("⏹️ Остановить анализ")           # Кнопка остановки анализа

    with col2:  # Правая колонка
        hcol2, linkcol2 = st.columns([12, 1])  # Заголовок + ссылка-якорь
        with hcol2:
            st.subheader("🔗 Сайты конкурентов")  # Подзаголовок секции
        with linkcol2:
            st.markdown('<div style="text-align:right; padding-top:10px;"><a href="#exp-comp-1">🔗</a></div>',
                        unsafe_allow_html=True)  # Ссылка-якорь к первому конкуренту внизу

        competitors = st.text_area("URL конкурентов (по одному на строку)", key="competitors",  # Многострочный ввод
                                   placeholder="https://конкурент1.com\nhttps://конкурент2.com",
                                   help="Сравните с конкурентами", height=100)

# «ЖИВАЯ» ЗОНА СТАТУСОВ
live_area_ph = st.empty()  # Плейсхолдер под динамические статусы анализа (всё рисуем внутри него)

# запуск/останов
if submitted:                                # Если нажали «Начать анализ»
    st.session_state.is_running = True       # Ставим флаг «идёт анализ»
    st.session_state.cancel = False          # Сбрасываем флаг отмены
    st.session_state.analyses_last = {}      # Чистим прошлые результаты
    st.session_state.urls_last = {}          # Чистим прошлые URL'ы
    st.session_state.shots_last = {}         # Чистим прошлые скриншоты
if stop_clicked:                             # Если нажали «Остановить анализ»
    st.session_state.cancel = True           # Ставим флаг отмены

# Анализ
if submitted or st.session_state.is_running:   # Если либо только что нажали, либо анализ ещё идёт
    if not st.session_state.get("main_url"):   # Если не указан основной URL —
        st.error("❌ Пожалуйста, укажите URL вашего сайта")  # Сообщаем
        st.session_state.is_running = False    # Снимаем флаг
        st.stop()                              # Останавливаем дальнейший рендер

    urls = {"Ваш сайт": st.session_state["main_url"].strip()}  # Начинаем словарь URL'ов с вашего сайта
    if st.session_state.get("competitors"):  # Если введены конкуренты —
        urls.update({f"Конкурент {i + 1}": u.strip()          # Добавляем каждого конкурента в словарь
                     for i, u in enumerate(st.session_state["competitors"].splitlines()) if u.strip()})

    if len(urls) > CONFIG["max_pages"]:  # Если превышаем лимит страниц —
        st.error(f"⚠️ Превышено максимальное количество страниц ({CONFIG['max_pages']})")  # Сообщаем
        st.session_state.is_running = False  # Снимаем флаг
        st.stop()                            # Останавливаем рендер

    with st.spinner(f"🔍 Анализируем {len(urls)} страниц..."):  # Плашка-спиннер на время анализа
        analyses: Dict[str, str] = {}           # Сюда будем складывать тексты анализа
        shots_all: Dict[str, List[BytesIO]] = {}  # Сюда — скриншоты по каждому сайту
        progress = st.progress(0)               # Прогресс-бар (0..1)

        with live_area_ph.container():          # Рисуем все статусы внутри общего контейнера (потом его очистим)
            for i, (name, url) in enumerate(urls.items()):  # Идём по всем URL'ам
                if st.session_state.cancel:   # Если пользователь прервал —
                    st.warning("Анализ остановлен пользователем.")  # Сообщаем
                    break                     # Выходим из цикла

                with st.status(f"Анализ: {name}", expanded=True):  # Статус-карточка для текущей страницы
                    st.write(f"🌐 Загрузка страницы: {url}")        # Пишем какой URL загружается

                    html = get_page_html(url)         # Загружаем HTML
                    if not html:
                        # Всегда записываем результат, чтобы конкурент появился в итогах
                        analyses[name] = "Не удалось загрузить страницу."
                        # Здесь НЕ дублируем st.error — get_page_html уже показал подробную ошибку.
                        st.write("⚠️ Страница не загружена")  # мягкое сообщение внутри статуса (опц.)
                    if st.session_state.cancel:       # Проверяем отмену после долгой операции
                        st.warning("Анализ остановлен пользователем.")
                        break

                    if html:                          # Если HTML получили —
                        info_slot = st.empty()        # Плейсхолдер для строки «Анализ контента...»
                        info_slot.write("🧠 Анализ контента...")  # Пишем «Анализ контента...»

                        res = analyze_page(html, name)  # Запускаем LLM анализ

                        if st.session_state.cancel:     # Снова проверяем отмену
                            st.warning("Анализ остановлен пользователем.")
                            break

                        if res:                         # Если пришёл ответ —
                            analyses[name] = res        # Сохраняем анализ
                            info_slot.success("✅ Анализ завершен!")  # Обновляем строку статуса

                            if include_screenshots:     # Если включены скриншоты —
                                try:
                                    shots = take_screenshots(url, max_shots=3)  # Делаем до 3 скринов
                                except Exception as e:
                                    st.warning(f"⚠️ Скриншоты не получились: {e}")  # Сообщаем, но не падаем
                                    shots = []  # Пустой список
                                shots_all[name] = shots  # Сохраняем скриншоты для этого сайта
                                st.info(f"📸 Скриншотов получено: {len(shots)}")  # Выводим количество

                            with st.expander("📄 Просмотреть результаты", expanded=False):  # Сворачиваемый блок с текстом анализа
                                st.markdown(res)  # Печатаем анализ
                        else:
                            info_slot.error("❌ Ошибка анализа страницы")  # Если LLM не ответила — сообщаем
                    else:
                        st.error("❌ Не удалось загрузить страницу")  # Если HTML не удалось получить

                progress.progress((i + 1) / len(urls))  # Обновляем прогресс-бар

        st.session_state.is_running = False  # По окончании цикла — сбрасываем флаг «идёт анализ»

        if st.session_state.cancel:          # Если был запрос на отмену
            st.session_state.cancel = False  # Сбрасываем флаг
            st.info("⏹️ Анализ прерван, введённые данные сохранены.")  # Сообщаем пользователю
            st.stop()  # Останавливаем рендер текущего кадра (UI останется с введёнными данными)

        if analyses:                         # Если что-то проанализировали —
            st.session_state.analyses_last = analyses   # Сохраняем результаты анализа
            st.session_state.urls_last = urls           # Сохраняем URL'ы
            st.session_state.shots_last = shots_all     # Сохраняем скриншоты
            live_area_ph.empty()                        # ЧИСТИМ «живую зону» со статусами (они исчезают сверху)

# Итоги
if st.session_state.analyses_last:  # Если есть результаты предыдущего анализа
    st.success("🎉 Все анализы успешно завершены!")  # Показываем «успех»

    st.subheader("📊 Результаты анализа")  # Заголовок секции с вкладками
    tabs = st.tabs(list(st.session_state.analyses_last.keys()))  # Создаём вкладки для каждой страницы
    for tab, (nm, txt) in zip(tabs, st.session_state.analyses_last.items()):  # Наполняем вкладки
        with tab:
            st.markdown(txt)  # Печатаем текст анализа

    # ✅ Защищаем генерацию PPTX, чтобы не ронять UI при любой ошибке
    try:
        pptx_io = create_ux_presentation_bytes(  # Формируем презентацию
            st.session_state.analyses_last,      # Тексты анализа
            st.session_state.shots_last,         # Скриншоты
            min_slides=3,                        # Минимум 3 слайда на сайт
            max_slides=5                         # Максимум 5 слайдов на сайт
        )
        st.download_button(                      # Кнопка скачивания отчёта
            "📥 Скачать отчет (PPTX)",           # Текст кнопки
            data=pptx_io,                        # Данные PPTX из памяти
            file_name=f"UX_Отчет_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",  # Имя файла
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"  # MIME тип PPTX
        )
    except Exception as e:                        # Ловим любые ошибки при сборке PPTX
        st.error(f"⚠️ Ошибка при формировании PPTX: {e}")  # Показываем ошибку, но не валим интерфейс

    # Свёрнутые «Анализ: …» внизу + якоря для наших ссылок
    st.markdown('<a id="exp-main"></a>', unsafe_allow_html=True)  # Якорь для «Анализируемый сайт»
    if "Ваш сайт" in st.session_state.analyses_last:              # Если анализ для вашего сайта есть —
        with st.expander("Анализ: Ваш сайт", expanded=False):     # Сворачиваемый блок
            url = st.session_state.urls_last.get("Ваш сайт", "")  # Достаём URL
            if url:
                st.write(f"🌐 Загрузка страницы: {url}")          # Печатаем URL
            st.success("✅ Анализ завершен!")                      # Плашка успешного завершения
            with st.expander("📄 Просмотреть результаты", expanded=False):  # Вложенный expander с текстом
                st.markdown(st.session_state.analyses_last["Ваш сайт"])     # Печать анализа

    idx_comp = 1  # Счётчик конкурентов для якорей
    for nm, txt in st.session_state.analyses_last.items():  # Идём по всем результатам
        if nm == "Ваш сайт":              # Пропускаем «Ваш сайт» — его вывели выше
            continue
        st.markdown(f'<a id="exp-comp-{idx_comp}"></a>', unsafe_allow_html=True)  # Якорь для конкретного конкурента
        with st.expander(f"Анализ: {nm}", expanded=False):  # Сворачиваемый блок по конкуренту
            url = st.session_state.urls_last.get(nm, "")    # URL конкурента
            if url:
                st.write(f"🌐 Загрузка страницы: {url}")    # Печатаем URL
            st.success("✅ Анализ завершен!")                # Плашка успеха
            with st.expander("📄 Просмотреть результаты", expanded=False):  # Вложенный expander с текстом
                st.markdown(txt)                            # Печать анализа
        idx_comp += 1                                       # Увеличиваем индекс якоря

# Инструкция — если ничего не запущено и нет результатов
if not st.session_state.is_running and not st.session_state.analyses_last:  # Если анализ не идёт и результатов нет
    st.info("""
    ℹ️ **Инструкция:**
        \nВедите в поля выше анализируемую
          страницу и страницы конкурентов, 
          относящиеся к тому же уровню - например, 
          сравниваем главную страницу с главными 
          или страницы выбора продукта между собой.
    """)  # Показываем инструкцию пользователю
